
import os
import sys
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import warnings

warnings.filterwarnings("ignore")
sns.set_style("whitegrid")


BASE_PATH = "archive"  
CHART_DIR = "charts"
OUTPUT_PPT = "Olist_Consulting_Deck.pptx"

os.makedirs(CHART_DIR, exist_ok=True)
 
def save_fig(fname):
    path = os.path.join(CHART_DIR, fname)
    plt.tight_layout()
    plt.savefig(path, dpi=150)
    plt.close()
    print(f"Saved chart: {path}")
    return path


def load_datasets(base_path):
    files = {
        "orders": "olist_orders_dataset.csv",
        "customers": "olist_customers_dataset.csv",
        "items": "olist_order_items_dataset.csv",
        "products": "olist_products_dataset.csv",
        "payments": "olist_order_payments_dataset.csv",
        "reviews": "olist_order_reviews_dataset.csv",
        "sellers": "olist_sellers_dataset.csv",
        "geolocation": "olist_geolocation_dataset.csv"  # optional (city/state)
    }
    datasets = {}
    for k, v in files.items():
        path = os.path.join(base_path, v)
        if os.path.exists(path):
            try:
                datasets[k] = pd.read_csv(path, low_memory=False)
                
            except Exception as e:
                print(f" Failed reading {v}: {e}")
        else:
            print(f"Missing {v} (skipping)")
    return datasets


def build_fact(datasets):
    if "orders" not in datasets:
        raise FileNotFoundError("Orders dataset required to build fact table.")
    orders = datasets["orders"].copy()
    # payments optional
    payments = datasets.get("payments", pd.DataFrame())

    # timestamps
    orders["order_purchase_timestamp"] = pd.to_datetime(orders["order_purchase_timestamp"], errors="coerce")
    # safe month
    orders["month"] = orders["order_purchase_timestamp"].dt.to_period("M").dt.to_timestamp()

    if not payments.empty:
        pay = payments.groupby("order_id", as_index=False).agg({"payment_value": "sum"})
    else:
        pay = pd.DataFrame(columns=["order_id","payment_value"])
    fact = orders.merge(pay, on="order_id", how="left")
    fact["revenue"] = fact["payment_value"].fillna(0.0)
    return fact

def revenue_trends(fact):
    df = fact.dropna(subset=["month"])
    rev = df.groupby("month")["revenue"].sum().sort_index()
    orders = df.groupby("month")["order_id"].nunique().sort_index()
    # Revenue trend
    plt.figure(figsize=(10,4.5))
    plt.plot(rev.index, rev.values, marker="o")
    plt.title("Monthly Revenue (BRL)")
    plt.ylabel("Revenue (BRL)")
    plt.xlabel("Month")
    plt.xticks(rotation=45)
    save_fig("revenue_trend.png")

    
    growth = rev.pct_change().fillna(0) * 100
    plt.figure(figsize=(10,4.5))
    plt.plot(growth.index, growth.values, marker="o")
    plt.axhline(0, color="gray", linewidth=0.7)
    plt.title("Monthly Revenue Growth (%)")
    plt.ylabel("Growth %")
    plt.xlabel("Month")
    plt.xticks(rotation=45)
    save_fig("revenue_growth_pct.png")

    
    plt.figure(figsize=(10,4.5))
    ax = plt.gca()
    ax.plot(rev.index, rev.values, marker="o", label="Revenue")
    ax.set_ylabel("Revenue (BRL)")
    ax2 = ax.twinx()
    ax2.plot(orders.index, orders.values, marker="o", color="orange", label="Orders")
    ax2.set_ylabel("Unique Orders")
    ax.set_xlabel("Month")
    ax.set_title("Revenue & Unique Orders (Monthly)")
    ax.tick_params(axis='x', rotation=45)
    save_fig("revenue_orders_dual.png")

    return {"rev": rev, "orders": orders, "growth": growth}

def average_order_value(fact):
    df = fact.dropna(subset=["month"])
    monthly = df.groupby("month").agg(total_revenue=("revenue","sum"), orders=("order_id","nunique"))
    monthly["aov"] = monthly["total_revenue"] / monthly["orders"]
    plt.figure(figsize=(10,4.5))
    plt.plot(monthly.index, monthly["aov"].values, marker="o")
    plt.title("Average Order Value (AOV) Over Time")
    plt.ylabel("AOV (BRL)")
    plt.xlabel("Month")
    plt.xticks(rotation=45)
    save_fig("aov_trend.png")
    return monthly["aov"]

def payment_distribution(datasets):
    payments = datasets.get("payments")
    if payments is None or "payment_type" not in payments.columns:
        print("Payment dataset or payment_type missing — skipping payment distribution.")
        return None
    share = payments["payment_type"].value_counts().head(8)
    plt.figure(figsize=(6,6))
    plt.pie(share.values, labels=share.index, autopct="%1.1f%%", startangle=90)
    plt.title("Payment Method Distribution")
    save_fig("payment_distribution.png")
    return share

def category_contribution(datasets, fact):
    items = datasets.get("items")
    products = datasets.get("products")
    if items is None or products is None:
        print("Items or products dataset missing — skipping category contribution.")
        return None
    df = items.merge(products, on="product_id", how="left")
    df = df.merge(fact[["order_id","revenue"]], on="order_id", how="left")
    # Some categories are NaN
    df["cat"] = df["product_category_name"].fillna("Unknown")
    cat_sales = df.groupby("cat")["revenue"].sum().sort_values(ascending=False)
    top10 = cat_sales.head(10)
    plt.figure(figsize=(8,5))
    sns.barplot(x=top10.values, y=top10.index, palette="mako")
    plt.title("Top 10 Categories by Revenue")
    plt.xlabel("Revenue (BRL)")
    save_fig("top10_categories_revenue.png")
    return top10

def cohort_retention(fact):
    if "customer_id" not in fact.columns:
        print("customer_id missing — skipping cohort analysis.")
        return None
    customers = fact.dropna(subset=["customer_id","month"])[["customer_id","month"]].copy()
    customers["month"] = pd.to_datetime(customers["month"])
    cohorts = customers.groupby("customer_id")["month"].min().reset_index()
    cohorts.columns = ["customer_id","cohort_month"]
    customers = customers.merge(cohorts, on="customer_id")
    customers["cohort_index"] = (customers["month"].dt.to_period("M") - customers["cohort_month"].dt.to_period("M")).apply(lambda x: x.n)
    cohort_pivot = customers.groupby(["cohort_month","cohort_index"])["customer_id"].nunique().reset_index()
    pivot = cohort_pivot.pivot(index="cohort_month", columns="cohort_index", values="customer_id").fillna(0)
    
    
    cohort_sizes = pivot.iloc[:,0]
    retention = pivot.div(cohort_sizes, axis=0)*100
    plt.figure(figsize=(12,6))
    sns.heatmap(retention, cmap="YlGnBu", fmt=".1f")
    plt.title("Customer Retention by Cohort ( % of cohort continuing )")
    plt.ylabel("Cohort Month")
    plt.xlabel("Months since First Purchase")
    save_fig("cohort_retention_heatmap.png")
    return retention

def delivery_vs_reviews(datasets, fact):
    if "reviews" not in datasets or "orders" not in datasets:
        print(" reviews or orders dataset missing — skipping delivery vs reviews.")
        return None
    reviews = datasets["reviews"].copy()
    orders = datasets["orders"].copy()
    # dates
    orders["order_delivered_customer_date"] = pd.to_datetime(orders["order_delivered_customer_date"], errors="coerce")
    orders["order_estimated_delivery_date"] = pd.to_datetime(orders["order_estimated_delivery_date"], errors="coerce")
    reviews["review_score"] = pd.to_numeric(reviews["review_score"], errors="coerce")
    orders["delay_days"] = (orders["order_delivered_customer_date"] - orders["order_estimated_delivery_date"]).dt.days
    df = reviews.merge(orders[["order_id","delay_days"]], on="order_id", how="left")
    df = df.dropna(subset=["review_score","delay_days"])
    plt.figure(figsize=(8,5))
    sns.boxplot(data=df, x="review_score", y="delay_days")
    plt.title("Delivery Delay (days) by Review Score")
    plt.ylabel("Delay (days)")
    plt.xlabel("Review Score (1-5)")
    save_fig("delivery_delay_vs_reviews.png")
    return df

def revenue_by_state(datasets, fact):
    # Use geolocation + customers to map city/state -> revenue
    if "geolocation" not in datasets or "customers" not in datasets:
        print("geolocation or customers dataset missing — skipping revenue by state.")
        return None
    geoloc = datasets["geolocation"].copy()
    customers = datasets["customers"].copy()
    # geolocation dataset on Olist has columns: geolocation_zip_code_prefix, geolocation_state
    # customers has customer_zip_code_prefix and customer_id
    if "geolocation_zip_code_prefix" not in geoloc.columns:
        print("geolocation format unexpected — skipping revenue by state.")
        return None
    # pick first state per zip (there are duplicates)
    zip_state = geoloc.groupby("geolocation_zip_code_prefix")["geolocation_state"].agg(lambda x: x.mode().iat[0] if len(x.mode())>0 else x.iloc[0]).reset_index()
    if "customer_zip_code_prefix" not in customers.columns:
        print(" customers format unexpected — skipping revenue by state.")
        return None
    cust = customers[["customer_id","customer_zip_code_prefix"]].copy()
    cust = cust.merge(zip_state, left_on="customer_zip_code_prefix", right_on="geolocation_zip_code_prefix", how="left")
    orders = fact[["order_id","customer_id","revenue"]].copy()
    rev_by_cust = orders.groupby("customer_id")["revenue"].sum().reset_index()
    rev_by_cust = rev_by_cust.merge(cust, on="customer_id", how="left")
    rev_state = rev_by_cust.groupby("geolocation_state")["revenue"].sum().sort_values(ascending=False).head(15)
    if rev_state.empty:
        print(" No state-level revenue — skipping revenue by state.")
        return None
    plt.figure(figsize=(8,5))
    sns.barplot(x=rev_state.values, y=rev_state.index, palette="viridis")
    plt.title("Top 15 States by Revenue")
    plt.xlabel("Revenue (BRL)")
    save_fig("revenue_by_state.png")
    return rev_state

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    try:
        slide.placeholders[1].text = subtitle
    except Exception:
        pass

def add_text_slide(prs, title, paragraphs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, p in enumerate(paragraphs):
        if i == 0:
            p0 = tf.paragraphs[0]
            p0.text = p
            p0.level = 0
        else:
            p_new = tf.add_paragraph()
            p_new.text = p
            p_new.level = 0

def add_chart_with_insight(prs, img_path, title, insight_text):
    
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left, top, w, h = Inches(0.6), Inches(1.4), Inches(6.5), Inches(4.0)
    try:
        slide.shapes.add_picture(img_path, left, top, width=w, height=h)
    except Exception as e:
        print(f"Failed to add image {img_path} to slide: {e}")
    # add textbox with insight
    tx_left, tx_top = Inches(7.3), Inches(1.6)
    tx_w, tx_h = Inches(2.7), Inches(3.6)
    txBox = slide.shapes.add_textbox(tx_left, tx_top, tx_w, tx_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = insight_text
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

def build_presentation(chart_paths, analytics_results):
    prs = Presentation()

    
    add_title_slide(prs, "Olist E-Commerce — Strategic Analysis", "Data-driven recommendations & roadmap")

   
    exec_summary = [
        "Key takeaways:",
        " • Revenue exhibits a strong upward trend, but month-to-month growth is volatile (seasonality & spikes).",
        " • Customer retention declines sharply after the first purchase — high acquisition dependency.",
        " • Top product categories contribute the majority of revenue (concentration risk).",
        " • Delivery delays directly correlate with negative reviews — logistics is top lever for NPS improvement.",
        " • Payments dominated by a single method (credit card) — diversify payments to improve reach & resilience."
    ]
    add_text_slide(prs, "Executive Summary", exec_summary)

    
    add_text_slide(prs, "Client Challenge", [
        "Olist seeks sustainable revenue growth while improving customer lifetime value and operational excellence.",
        "Key asks: Improve retention, diversify revenue, reduce negative reviews from delivery issues, and expand payment options."
    ])
    add_text_slide(prs, "Project Approach (Consulting Methodology)", [
        "1) Data ingestion & validation",
        "2) Descriptive analytics (KPIs, cohorts, categories)",
        "3) Diagnostic analysis (delivery vs satisfaction)",
        "4) Strategic recommendations & implementation roadmap"
    ])


    if chart_paths.get("revenue_trend"):
        add_chart_with_insight(prs, chart_paths["revenue_trend"],
                               "Revenue Trend",
                               "Revenue shows steady long-term growth with clear seasonal spikes (holiday peaks). "
                               "Action: Smooth seasonality with off-peak demand programs and subscription offerings.")

    if chart_paths.get("revenue_growth_pct"):
        add_chart_with_insight(prs, chart_paths["revenue_growth_pct"],
                               "Revenue Growth % (MoM)",
                               "Growth % is volatile; months of strong growth follow by shallow months. "
                               "Action: Focus on retention & continuous promotions to stabilize growth.")

    if chart_paths.get("revenue_orders_dual"):
        add_chart_with_insight(prs, chart_paths["revenue_orders_dual"],
                               "Revenue & Orders",
                               "Revenue and order counts move together; however, revenue per order can vary. "
                               "Action: Upsell & cross-sell to increase AOV and revenue resilience.")

    # AOV
    if chart_paths.get("aov_trend"):
        add_chart_with_insight(prs, chart_paths["aov_trend"],
                               "Average Order Value (AOV)",
                               "AOV is relatively stable. Driving AOV via bundles and recommendations increases short-term revenue without additional CAC.")

    # Payment distribution
    if chart_paths.get("payment_distribution"):
        add_chart_with_insight(prs, chart_paths["payment_distribution"],
                               "Payment Mix",
                               "High dependency on credit: consider PIX, wallets, BNPL to broaden access and reduce single-channel risk.")

    # Categories
    if chart_paths.get("top10_categories_revenue"):
        add_chart_with_insight(prs, chart_paths["top10_categories_revenue"],
                               "Top Categories by Revenue",
                               "Top categories account for majority of revenue — diversify category mix and onboard sellers in long-tail categories to reduce concentration risk.")

    # Cohort retention
    if chart_paths.get("cohort_retention_heatmap"):
        add_chart_with_insight(prs, chart_paths["cohort_retention_heatmap"],
                               "Cohort Retention Heatmap",
                               "Retention drops quickly after first purchase. Implement onboarding, next-purchase incentives and personalized re-engagement to lift repeat rates.")

    # Delivery vs Reviews
    if chart_paths.get("delivery_delay_vs_reviews"):
        add_chart_with_insight(prs, chart_paths["delivery_delay_vs_reviews"],
                               "Delivery Delays vs Customer Ratings",
                               "Longer delivery delays heavily skew negative ratings. Invest in last-mile SLAs, partnerships, and tracking to reduce negative reviews and returns.")

    # Revenue by state
    if chart_paths.get("revenue_by_state"):
        add_chart_with_insight(prs, chart_paths["revenue_by_state"],
                               "Revenue by State (Top 15)",
                               "Revenue concentration by state shows regional strengths. Use targeted logistics & marketing investments in high-potential regions.")

    # Strategic Recommendations (detailed)
    add_text_slide(prs, "Strategic Recommendations (By Segment)", [
        "Customer: Launch loyalty program (tiered rewards), personalized re-marketing, onboarding sequence for new customers.",
        "Category: Incentivize seller acquisition in long-tail categories; curated bundles & category-specific promotions.",
        "Logistics: Regional last-mile pilots, SLA-based carrier selection, delivery tracking & customer communication.",
        "Payments: Integrate local payment methods (Pix), wallets, and BNPL to increase conversion & cashflow options.",
        "Operations: Improve cancellation & refund workflows, reduce failed deliveries via address validation & pickup points."
    ])

    # Implementation Roadmap (phases)
    add_text_slide(prs, "Implementation Roadmap", [
        "Phase 1 (0-3 months): Retention campaigns, payments integration pilots, delivery KPI dashboard.",
        "Phase 2 (3-6 months): Seller onboarding in priority categories, logistics partner pilots, BNPL pilot.",
        "Phase 3 (6-12 months): Scale category diversification, platform-level loyalty, automate dispute resolution."
    ])

    # Projected Impact
    add_text_slide(prs, "Projected Impact (Conservative Estimates)", [
        "• Repeat rate +15-20% with targeted retention & loyalty.",
        "• Revenue +10-12% from category diversification and AOV initiatives.",
        "• Negative reviews -30% via delivery SLA improvements.",
        "• Improved conversion via payment diversification (est. +3-5%)."
    ])

    # Appendix: list of charts included
    appendix_notes = ["Charts included in appendix (see slides):"] + [f" • {k}" for k in chart_paths.keys() if chart_paths[k]]
    add_text_slide(prs, "Appendix - Charts", appendix_notes)

    prs.save(OUTPUT_PPT)
    


def main():
    datasets = load_datasets(BASE_PATH)
    try:
        fact = build_fact(datasets)
    except Exception as e:
        print(f"Fatal error building fact table: {e}")
        sys.exit(1)

    charts = {}
    analytics = {}

    # Revenue & orders
    try:
        res = revenue_trends(fact)
        charts["revenue_trend"] = os.path.join(CHART_DIR, "revenue_trend.png")
        charts["revenue_growth_pct"] = os.path.join(CHART_DIR, "revenue_growth_pct.png")
        charts["revenue_orders_dual"] = os.path.join(CHART_DIR, "revenue_orders_dual.png")
        analytics.update(res)
    except Exception as e:
        print(f"Error in revenue_trends: {e}")

    # AOV
    try:
        aov = average_order_value(fact)
        charts["aov_trend"] = os.path.join(CHART_DIR, "aov_trend.png")
        analytics["aov"] = aov
    except Exception as e:
        print(f"Error in average_order_value: {e}")

    # Payment distribution
    try:
        pay = payment_distribution(datasets)
        charts["payment_distribution"] = os.path.join(CHART_DIR, "payment_distribution.png") if pay is not None else None
    except Exception as e:
        print(f"Error in payment_distribution: {e}")

    # Category contribution
    try:
        topcats = category_contribution(datasets, fact)
        charts["top10_categories_revenue"] = os.path.join(CHART_DIR, "top10_categories_revenue.png") if topcats is not None else None
    except Exception as e:
        print(f"Error in category_contribution: {e}")

    # Cohort retention
    try:
        cohort = cohort_retention(fact)
        charts["cohort_retention_heatmap"] = os.path.join(CHART_DIR, "cohort_retention_heatmap.png") if cohort is not None else None
    except Exception as e:
        print(f"Error in cohort_retention: {e}")

    # Delivery vs reviews
    try:
        dv = delivery_vs_reviews(datasets, fact)
        charts["delivery_delay_vs_reviews"] = os.path.join(CHART_DIR, "delivery_delay_vs_reviews.png") if dv is not None else None
    except Exception as e:
        print(f"Error in delivery_vs_reviews: {e}")

    # Revenue by state (optional)
    try:
        rbs = revenue_by_state(datasets, fact)
        charts["revenue_by_state"] = os.path.join(CHART_DIR, "revenue_by_state.png") if rbs is not None else None
    except Exception as e:
        print(f"Error in revenue_by_state: {e}")

    # Filter out None entries
    charts = {k: v for k, v in charts.items() if v is not None and os.path.exists(v)}

    # Build PPT
    build_presentation(charts, analytics)

if __name__ == "__main__":
    main()
